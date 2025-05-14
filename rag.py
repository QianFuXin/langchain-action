import os
import time
import uuid
import json
import logging
import requests

import pandas as pd
from flask import Flask, request, jsonify
from flask_cors import CORS

from elasticsearch import Elasticsearch
from elasticsearch.helpers import bulk

from langchain_chroma import Chroma
from langchain_core.prompts import PromptTemplate
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter

from docx import Document
from pptx import Presentation

from config import api_key
from utils import embeddings, model

# === 日志配置 ===
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

# === Elasticsearch 初始化 ===
try:
    es = Elasticsearch("http://127.0.0.1:9200")
    if not es.ping():
        raise Exception("无法连接本地 Elasticsearch")
except Exception as e:
    logging.error(f"Elasticsearch 初始化失败：{e}")
    raise

# === Flask 应用配置 ===
app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = "uploads"
VECTOR_DATA = "./chroma_data"
MAX_CONTENT_LENGTH = 20 * 1024 * 1024
ALLOWED_EXTENSIONS = {'pdf', 'xls', 'xlsx', 'docx', 'pptx', 'txt'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# === 通用文本分割器 ===
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, add_start_index=True)

# === Prompt 模板 ===
prompt = PromptTemplate.from_template(
    """
    你是一名专注于问答任务的助手。请根据以下提供的上下文内容回答问题。
    如果无法从上下文中获得答案，请明确表示你不知道。
    回答应简洁明了，最多不超过三句话。

    问题： {question}

    上下文： {context}

    回答：
    """
)


# === 工具函数 ===

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[-1].lower() in ALLOWED_EXTENSIONS


def split_text(content: str):
    return text_splitter.split_text(content)


def get_excel_data(file_path):
    df = pd.read_excel(file_path, dtype=str)
    return df.apply(lambda row: "|".join([f"{col}:{row[col]}" for col in df.columns]), axis=1).tolist()


def get_pdf_data(file_path):
    loader = PyPDFLoader(file_path)
    docs = loader.load()
    return split_text("".join([doc.page_content for doc in docs]))


def get_word_data(file_path):
    doc = Document(file_path)
    full_text = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
    return split_text(full_text)


def get_ppt_data(file_path):
    prs = Presentation(file_path)
    full_text = "\n".join(
        [shape.text.strip() for slide in prs.slides for shape in slide.shapes if
         hasattr(shape, "text") and shape.text.strip()]
    )
    return split_text(full_text)


def get_txt_data(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()
    return split_text(content)


def batch_data(data, batch_size=60):
    for i in range(0, len(data), batch_size):
        yield data[i:i + batch_size]


def rerank_documents(query, documents, model_name="BAAI/bge-reranker-v2-m3", top_n=4):
    url = "https://api.siliconflow.cn/v1/rerank"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }
    payload = {
        "model": model_name,
        "query": query,
        "documents": documents,
        "top_n": top_n,
        "return_documents": False,
        "max_chunks_per_doc": 1024,
        "overlap_tokens": 80
    }
    try:
        response = requests.post(url, headers=headers, data=json.dumps(payload))
        response.raise_for_status()
        results = response.json().get("results", [])
        sorted_results = sorted(results, key=lambda x: x["relevance_score"], reverse=True)[:top_n]
        return [documents[item["index"]] for item in sorted_results]
    except Exception as e:
        logging.error(f"Rerank API 调用失败: {e}")
        raise


# === 上传接口 ===
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"code": 400, "msg": "未上传文件"})

    file = request.files['file']
    if not file or not allowed_file(file.filename):
        return jsonify({"code": 400, "msg": "不支持的文件类型"})

    file_ext = file.filename.rsplit('.', 1)[-1].lower()
    unique_filename = f"{uuid.uuid4().hex}.{file_ext}"
    file_path = os.path.join(UPLOAD_FOLDER, unique_filename)
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    file.save(file_path)

    dataset_id = str(int(time.time()))
    try:
        loader_map = {
            'xls': get_excel_data,
            'xlsx': get_excel_data,
            'pdf': get_pdf_data,
            'docx': get_word_data,
            'pptx': get_ppt_data,
            'txt': get_txt_data
        }
        ai_data = loader_map[file_ext](file_path)

        # 写入向量数据库
        vector_store = Chroma(
            collection_name=dataset_id,
            embedding_function=embeddings,
            persist_directory=VECTOR_DATA,
        )
        for batch in batch_data(ai_data):
            vector_store.add_texts(batch)

        # 写入 Elasticsearch
        if not es.indices.exists(index=dataset_id):
            es.indices.create(index=dataset_id)
        bulk(es, [{"_index": dataset_id, "_source": {"content": i}} for i in ai_data])

        return jsonify({"code": 200, "data": {"dataset_id": dataset_id}, "msg": "上传并索引成功"})

    except Exception as e:
        logging.error(f"[{dataset_id}] 文件处理失败：{e}")
        return jsonify({"code": 500, "msg": "服务内部错误", "error": str(e)})


# === 查询接口 ===
@app.route('/query', methods=['POST'])
def query_answer():
    data = request.json
    if not data or 'dataset_id' not in data or 'query' not in data:
        return jsonify({"code": 400, "msg": "缺少 dataset_id 或 query"})

    dataset_id = data['dataset_id']
    query = data['query']
    try:
        # 向量查询
        vector_store = Chroma(
            collection_name=dataset_id,
            embedding_function=embeddings,
            persist_directory=VECTOR_DATA,
        )
        retrieved_data = [doc.page_content for doc in vector_store.similarity_search(query, k=4)]

        # 关键词查询
        search_query = {
            "query": {
                "multi_match": {
                    "query": query,
                    "fields": ["content"]
                }
            },
            "_source": ["content"],
            "size": 4
        }
        search_data = [hit["_source"]["content"] for hit in
                       es.search(index=dataset_id, body=search_query)["hits"]["hits"]]

        # 结果合并 & rerank
        all_data = list(set(retrieved_data + search_data))
        re_ranked = rerank_documents(query, all_data, top_n=4)
        context = "\n".join(re_ranked)

        # 构建 prompt 并调用模型
        prompt_msg = prompt.invoke({"question": query, "context": context})
        response = model.invoke(prompt_msg)

        return jsonify({"code": 200, "answer": response.content})

    except Exception as e:
        logging.error(f"[{dataset_id}] 查询失败：{e}")
        return jsonify({"code": 500, "msg": "查询失败", "error": str(e)})


# === 启动服务 ===
if __name__ == '__main__':
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    os.makedirs(VECTOR_DATA, exist_ok=True)
    app.run(host='0.0.0.0', port=5001, threaded=True)
